"""Proposal pitch — clean Beamer style, plain language, lots of breathing room."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LAVENDER    = RGBColor(0xC8, 0xC8, 0xE8)
LAVEN_LIGHT = RGBColor(0xEC, 0xEC, 0xF6)
NAVY        = RGBColor(0x1E, 0x1E, 0x6E)
INDIGO      = RGBColor(0x32, 0x32, 0x96)
GRAY        = RGBColor(0x66, 0x66, 0x66)
MID_GRAY    = RGBColor(0x99, 0x99, 0x99)
FONT        = "Palatino Linotype"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide):
    s = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = WHITE; s.line.fill.background()

def box(slide, x, y, w, h, fill, border=False):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = LAVENDER; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def t(slide, text, x, y, w, h, size, bold=False, color=NAVY,
      align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color; r.font.name = FONT
    return tb

def header(slide, title, n, total=6):
    bg(slide)
    box(slide, 0, 0, 13.33, 0.58, LAVENDER)
    t(slide, title, 0.3, 0.09, 11.5, 0.44, 22, color=NAVY)
    t(slide, f"{n} / {total}", 12.0, 0.12, 1.1, 0.35, 12, color=NAVY, align=PP_ALIGN.RIGHT)

def badge(slide, num, x, y):
    d = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(0.42), Inches(0.42))
    d.fill.solid(); d.fill.fore_color.rgb = INDIGO; d.line.fill.background()
    tf = d.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = str(num); r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = FONT

def callout(slide, title, body, x, y, w):
    """Simple block callout — header bar + light body."""
    box(slide, x, y, w, 0.42, INDIGO)
    t(slide, title, x + 0.15, y + 0.06, w - 0.3, 0.34, 14, color=WHITE)
    lines = body if isinstance(body, list) else [body]
    bh = 0.18 + len(lines) * 0.46
    box(slide, x, y + 0.42, w, bh, LAVEN_LIGHT)
    for i, line in enumerate(lines):
        t(slide, line, x + 0.18, y + 0.5 + i * 0.46, w - 0.36, 0.42, 14, color=NAVY)
    return y + 0.42 + bh


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1)

box(s1, 1.0, 1.6, 11.33, 1.7, LAVENDER)
t(s1, "Snake RL", 1.2, 1.72, 10.9, 0.85, 40, color=NAVY, align=PP_ALIGN.CENTER)
t(s1, "Teaching a computer to play Snake using Reinforcement Learning",
  1.2, 2.52, 10.9, 0.6, 17, color=NAVY, align=PP_ALIGN.CENTER)

t(s1, "Carlos Lavernia",
  1.2, 3.8, 10.9, 0.48, 20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
t(s1, "Advanced Machine Learning  ·  Capstone Proposal",
  1.2, 4.35, 10.9, 0.4, 15, color=GRAY, align=PP_ALIGN.CENTER)
t(s1, "April 2026",
  1.2, 4.85, 10.9, 0.38, 14, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

t(s1, "1 / 6", 12.0, 7.1, 1.1, 0.3, 12, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What & Why
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
header(s2, "What is This Project?", 2)

t(s2, "The Goal", 0.5, 0.78, 12.3, 0.42, 17, bold=True, color=INDIGO)
t(s2, "Build an AI agent that learns to play Snake entirely on its own — "
      "no hand-coded rules, no human guidance. Just trial and error.",
  0.5, 1.2, 12.3, 0.7, 16, color=NAVY)

box(s2, 0.5, 2.05, 12.33, 0.03, LAVENDER)

t(s2, "Why is this hard?", 0.5, 2.22, 12.3, 0.42, 17, bold=True, color=INDIGO)

items = [
    ("The agent can't see the future.",
     "It must learn from experience that cutting off its own path leads to death later."),
    ("Food gets harder to reach.",
     "As the snake grows longer, finding a safe route to the next apple gets much harder."),
    ("The win condition is extreme.",
     "Filling the entire 10×10 board means collecting 97 apples in a row without dying once."),
]
for i, (bold_txt, plain_txt) in enumerate(items):
    y = 2.75 + i * 1.3
    badge(s2, i + 1, 0.5, y + 0.04)
    t(s2, bold_txt,  1.08, y,        11.5, 0.42, 15, bold=True,  color=NAVY)
    t(s2, plain_txt, 1.08, y + 0.42, 11.5, 0.6,  14, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — How It Learns (plain RL explanation)
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
header(s3, "How Does the Agent Learn?", 3)

t(s3, "Reinforcement Learning in plain terms", 0.5, 0.78, 12.3, 0.42, 17, bold=True, color=INDIGO)
t(s3, "The agent plays thousands of games. After each action it receives a reward signal — "
      "positive for eating food, negative for dying. Over time it figures out which moves lead to better outcomes.",
  0.5, 1.22, 12.3, 0.82, 15, color=NAVY)

box(s3, 0.5, 2.18, 12.33, 0.03, LAVENDER)

# Three-column loop diagram (simplified)
steps = [
    ("1  Observe",   "The agent sees the\ngame board as a\n7-layer grid image."),
    ("2  Act",       "It picks a direction\n(up, down, left, right)\nbased on what it learned."),
    ("3  Learn",     "It receives a reward\nand updates its strategy\nto do better next time."),
]
for i, (title, body) in enumerate(steps):
    cx = 0.5 + i * 4.28
    box(s3, cx, 2.38, 3.9, 0.44, INDIGO)
    t(s3, title, cx + 0.15, 2.44, 3.6, 0.36, 15, bold=True, color=WHITE)
    box(s3, cx, 2.82, 3.9, 1.5, LAVEN_LIGHT)
    t(s3, body, cx + 0.18, 2.92, 3.55, 1.3, 14, color=NAVY)

# Arrow labels between boxes
t(s3, "→", 4.42, 3.3, 0.4, 0.4, 22, color=INDIGO, align=PP_ALIGN.CENTER)
t(s3, "→", 8.7,  3.3, 0.4, 0.4, 22, color=INDIGO, align=PP_ALIGN.CENTER)

box(s3, 0.5, 4.5, 12.33, 0.03, LAVENDER)
t(s3, "Reward signals used in this project", 0.5, 4.65, 12.3, 0.38, 15, bold=True, color=INDIGO)

rewards = [
    "Eating food: +10 to +50 points  (bigger reward as the snake gets longer)",
    "Dying (wall or self): −10 points",
    "Running out of time without eating: −15 points  (discourages circling)",
    "Filling the entire board: +500 points  (the jackpot)",
]
for i, r in enumerate(rewards):
    t(s3, f"• {r}", 0.65, 5.1 + i * 0.46, 12.1, 0.42, 13, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Technical Approach (accessible)
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
header(s4, "Technical Approach", 4)

t(s4, "The four main components of the system", 0.5, 0.75, 12.3, 0.42, 16, bold=True, color=INDIGO)

components = [
    ("A neural network that sees the board",
     "The agent uses a convolutional neural network (like those used in image recognition) "
     "to read the game grid and decide which direction to move."),
    ("A memory bank of past experiences",
     "The agent stores its last 100,000 moves and learns from random samples of them — "
     "this prevents it from overlearning recent events and forgetting older lessons."),
    ("A smarter learning algorithm",
     "Instead of crediting only the last move for a good outcome, the agent looks back "
     "5 steps. This helps it understand that planning ahead is what leads to success."),
    ("Structured exploration",
     "Early on, the agent moves randomly to explore. Over 2,000 games it gradually shifts "
     "to trusting what it has learned rather than guessing."),
]
for i, (title, body) in enumerate(components):
    y = 1.3 + i * 1.44
    badge(s4, i + 1, 0.5, y + 0.06)
    t(s4, title, 1.1, y,        12.0, 0.44, 15, bold=True, color=NAVY)
    t(s4, body,  1.1, y + 0.44, 12.0, 0.75, 13, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Expected Challenges
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
header(s5, "Expected Challenges", 5)

t(s5, "What will be hard to solve", 0.5, 0.75, 12.3, 0.42, 16, bold=True, color=INDIGO)

challenges = [
    ("The agent can get stuck in loops.",
     "Without the right incentive, it learns that circling indefinitely keeps it alive — "
     "even if it never eats. The timeout penalty is designed to break this habit."),
    ("Planning far into the future is really hard.",
     "The agent needs to think 20+ moves ahead to avoid trapping itself. "
     "Most RL agents struggle with this; it's the core challenge of the project."),
    ("Winning is extremely rare.",
     "The agent might become very good (score 60–80) without ever filling the full board. "
     "Getting from 'good' to 'perfect' requires a different level of strategic play."),
]
for i, (title, body) in enumerate(challenges):
    y = 1.3 + i * 1.8
    badge(s5, i + 1, 0.5, y + 0.06)
    t(s5, title, 1.1, y,        12.0, 0.44, 15, bold=True, color=NAVY)
    t(s5, body,  1.1, y + 0.44, 12.0, 0.85, 14, color=GRAY)
    if i < 2:
        box(s5, 0.5, y + 1.6, 12.33, 0.02, LAVENDER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Timeline
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
header(s6, "Timeline & What Success Looks Like", 6)

milestones = [
    ("Apr 13–15", "Proposal",      "This presentation"),
    ("Apr 14–20", "Build & Train", "Agent training,\nperformance tuning"),
    ("Apr 21–23", "Checkpoint",    "Demo + training\ncurves ready"),
    ("Apr 24–26", "Write-up",      "Report + analysis\ncompleted"),
    ("Apr 27–29", "Final Pres.",   "Live demo +\n10-min talk"),
]

box(s6, 0.5, 3.4, 12.33, 0.05, LAVENDER)

for i, (date, title, desc) in enumerate(milestones):
    cx = 0.38 + i * 2.5
    dot = s6.shapes.add_shape(9, Inches(cx + 0.8), Inches(3.22), Inches(0.42), Inches(0.42))
    dot.fill.solid(); dot.fill.fore_color.rgb = INDIGO; dot.line.fill.background()

    t(s6, date,  cx + 0.08, 1.35, 2.35, 0.38, 12, bold=True,  color=INDIGO, align=PP_ALIGN.CENTER)
    t(s6, title, cx + 0.08, 1.78, 2.35, 0.44, 14, bold=True,  color=NAVY,   align=PP_ALIGN.CENTER)
    box(s6, cx + 0.08, 2.28, 2.35, 0.82, LAVEN_LIGHT, border=True)
    t(s6, desc,  cx + 0.18, 2.34, 2.15, 0.72, 12,             color=GRAY,   align=PP_ALIGN.CENTER)

t(s6, "What success looks like", 0.5, 4.05, 12.3, 0.42, 16, bold=True, color=INDIGO)

goals = [
    "The agent consistently scores 50+ apples per game (out of 97 possible)",
    "The agent completes at least one full board — filling every cell",
    "Clear improvement visible in the score chart from episode 1 to episode 6,000",
]
for i, g in enumerate(goals):
    badge(s6, i + 1, 0.5, 4.6 + i * 0.72)
    t(s6, g, 1.1, 4.6 + i * 0.72, 12.0, 0.6, 14, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
out = "/Users/carlitoslavernia/Desktop/Snake_RL_Proposal.pptx"
prs.save(out)
print(f"Saved: {out}")
